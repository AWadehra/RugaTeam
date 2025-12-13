from pathlib import Path
import json
import hashlib
from datetime import datetime
from uuid import uuid4
import chromadb
import pymupdf4llm

# import pymupdf
from pptx import Presentation
from sentence_transformers import SentenceTransformer

# model = SentenceTransformer("all-MiniLM-L6-v2")
client = chromadb.PersistentClient(path="./P2/chroma_db")
collection = client.get_or_create_collection(name="documents")


def extract_pdf(path):
    # doc = pymupdf.open(path)
    # return " ".join([page.get_text() for page in doc])

    return pymupdf4llm.to_markdown(path)


def extract_pptx(path):
    prs = Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return " ".join(text)


def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> list[str]:
    # Try to break at sentence ends instead of mid-word
    chunks = []
    start = 0
    while start < len(text):
        end = min(start + chunk_size, len(text))
        # Look for last period/newline within chunk to break cleanly
        if end < len(text):
            for sep in ['\n\n', '\n', '. ', ' ']:
                last_break = text.rfind(sep, start, end)
                if last_break > start + chunk_size // 2:
                    end = last_break + len(sep)
                    break
        chunks.append(text[start:end].strip())
        start = end - overlap if end < len(text) else end
    return [c for c in chunks if c]  # filter empty

def llm_analysis(file, text):
    #"suggested_filename"
    #"categories":
    #"authors"
    #"topics"
    #"tags"
    #"summary"
    #"glossary_terms"
    
    prompt = f"""You are a document analyzer for a clinical epidemiology research department. Analyze this document and return a JSON object with the following fields.

            CATEGORIES (choose one with the highest confidence):
            - Education
            - Capita Selecta
            - Research Meeting
            - World Headlines
            - Miscellaneous

            FILE NAMING FORMAT: YYYY-MM-DD__Theme__Speaker__CoreSubject{file.suffix}
            - YYYY-MM-DD: Extract date from content, or use "UNKNOWN" if not found
            - Theme: One of [Education, Research, Methods, Clinical, Statistics]
            - Speaker: Author's last name, or "Unknown" if not found
            - CoreSubject: 2-4 word summary of main topic, normal cased with capital first letter without space (e.g., SurvivalAnalysis)

            Return this exact JSON structure (all string values, lists as comma-separated strings):
            {{
                "suggested_filename": "YYYY-MM-DD__Theme__Speaker__CoreSubject{file.suffix}",
                "categories": "category",
                "authors": "Author Name 1,Author Name 2",
                "topics": "topic1,topic2,topic3",
                "tags": "tag1,tag2,tag3",
                "summary": "2-3 sentence summary of the document",
                "glossary_terms": "term1,term2,term3"
            }}

            RULES:
            - authors: Full names, comma-separated. Use "Unknown" if not found.
            - topics: Main subjects/concepts covered (3-6 items)
            - tags: Broader categorization keywords like "methods", "statistics", "clinical research", "tutorial" (3-5 items)
            - glossary_terms: Specialized/technical terms that may need definition (medical, statistical jargon)
            - summary: Brief, informative summary of the document content

            DOCUMENT CONTENT:
            {text[:6000]}

            Return ONLY valid JSON, no additional text."""

    response = llm_call(prompt)
    
    return json.loads(response)
    
   


def create_metadata(file, text, chunk_no, analysis, file_id, content_hash):
    metadata = {
        "file_id": file_id,
        "original_path": str(file),
        "file_type":file.suffix[1:],
        "content_hash":content_hash,
        "title":file.stem,
        "chunk_no":chunk_no,
        #"suggested_filename"
        #"categories"
        "creation_date": datetime.fromtimestamp(file.stat().st_ctime).isoformat(),
        "last_modified_date": datetime.fromtimestamp(file.stat().st_mtime).isoformat(),
        #"analysis_date" -> No idea how to add this atm
        #"authors"
        #"topics"
        #"tags"
        #"summary"
        #"glossary_terms"
        "possible_duplicate":False,
        "reviewed_by_human":False,
        "llm_model":"Opus4.5",
        "extracted_at":datetime.now().isoformat(),
            
        **analysis
    }
    return metadata

# Indexing the folder
data_folder = Path("./data")

for file in data_folder.iterdir():
    if file.suffix == ".pdf":
        text = extract_pdf(file)
    elif file.suffix == ".pptx":
        text = extract_pptx(file)
    else:
        continue

    chunks = chunk_text(text)
    analysis = llm_analysis(file, text)
    file_id = str(uuid4())
    content_hash = hashlib.sha256(text.encode()).hexdigest()
    for i, chunk in enumerate(chunks):
        collection.add(
            ids=[f"{file.name}_chunk_{i}"],
            documents=[chunk],
            metadatas=[create_metadata(file, chunk, i, analysis, file_id, content_hash)],
        )
        print(f"Indexed: {file.name}")

# Result part
results = collection.query(query_texts=["survival analysis"], n_results=3)

print(results)
